from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

MAROON = RGBColor(0x50, 0x00, 0x00)
DARK = RGBColor(0x22, 0x22, 0x22)
GREY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_title_bar(slide, text: str) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.9))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = MAROON
    tf = bar.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.18)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = WHITE


def add_section(slide, left, top, width, height, heading: str, bullets: list[str]) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True

    head_p = tf.paragraphs[0]
    head_run = head_p.add_run()
    head_run.text = heading
    head_run.font.size = Pt(20)
    head_run.font.bold = True
    head_run.font.color.rgb = MAROON
    head_p.space_after = Pt(6)

    for bullet in bullets:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"• {bullet}"
        run.font.size = Pt(14)
        run.font.color.rgb = DARK
        p.space_after = Pt(4)


def add_system_card(slide, left, top, width, height, name: str, uni: str, desc: str, cite: str) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_top = Inches(0.02)

    p1 = tf.paragraphs[0]
    r_name = p1.add_run()
    r_name.text = name
    r_name.font.size = Pt(13)
    r_name.font.bold = True
    r_name.font.color.rgb = MAROON
    r_uni = p1.add_run()
    r_uni.text = f"  —  {uni}"
    r_uni.font.size = Pt(12)
    r_uni.font.bold = True
    r_uni.font.color.rgb = DARK
    p1.space_after = Pt(2)

    p2 = tf.add_paragraph()
    r_desc = p2.add_run()
    r_desc.text = desc
    r_desc.font.size = Pt(11)
    r_desc.font.color.rgb = DARK
    p2.space_after = Pt(2)

    p3 = tf.add_paragraph()
    r_cite = p3.add_run()
    r_cite.text = cite
    r_cite.font.size = Pt(9)
    r_cite.font.italic = True
    r_cite.font.color.rgb = GREY


def add_category_header(slide, left, top, width, text: str) -> None:
    box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = MAROON


def add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(10)
    run.font.italic = True
    run.font.color.rgb = GREY


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1 — Problem Statement
    s0 = prs.slides.add_slide(blank)
    add_title_bar(s0, "TamuBot — Problem Statement")

    add_section(
        s0,
        Inches(0.5),
        Inches(1.1),
        Inches(6.2),
        Inches(2.8),
        "Choosing Courses Takes Time",
        [
            "Searching for, reading, and understanding each syllabus is slow",
            "Comparing courses to find the best personal fit takes effort",
            "Easy to miss good courses you didn't know to look for",
        ],
    )

    add_section(
        s0,
        Inches(6.9),
        Inches(1.1),
        Inches(6.0),
        Inches(2.8),
        "Planning Around Constraints",
        [
            "Fitting courses to your specific program requirements",
            "Aligning class schedules across multiple courses",
            "Tracking prerequisites and registration windows term-to-term",
        ],
    )

    add_section(
        s0,
        Inches(0.5),
        Inches(4.0),
        Inches(6.2),
        Inches(2.9),
        "Advisor Time Is Limited",
        [
            "Advisors are a great resource, but appointments are short and scheduled",
            "Students get more out of meetings when they come prepared",
            "Self-research is the bottleneck before that conversation can start",
        ],
    )

    add_section(
        s0,
        Inches(6.9),
        Inches(4.0),
        Inches(6.0),
        Inches(2.9),
        "AI Without Grounding Falls Short",
        [
            "LLMs make broad knowledge easy to access",
            "Generic LLM answers are often not grounded in real documents",
            "Hallucinations make answers unreliable for high-stakes decisions",
            "Students need answers tied to actual TAMU syllabi they can verify",
        ],
    )

    add_footer(s0, "Texas A&M graduate course/policy assistant — RAG-grounded, citation-first")

    # Slide 2 — Current Status
    s1 = prs.slides.add_slide(blank)
    add_title_bar(s1, "TamuBot — Current Status")

    add_section(
        s1,
        Inches(0.5),
        Inches(1.1),
        Inches(6.2),
        Inches(2.8),
        "RAG Pipeline (live)",
        [
            "3-stage LangGraph: Router → Retrieval → Generator",
            "Gemini 2.5 Flash for routing, generation, PDF parsing",
            "Hybrid retrieval: vector + BM25 (RRF) + Voyage rerank-2",
            "8-function routing matrix (metadata/hybrid/semantic/compare)",
            "Recursive discovery: from an anchor course, find related ones",
            "Conversation memory via mem0 + SQLite checkpoints",
        ],
    )

    add_section(
        s1,
        Inches(6.9),
        Inches(1.1),
        Inches(6.0),
        Inches(2.8),
        "Data & Ingestion (7-step pipeline)",
        [
            "1. Scrape — Howdy Portal + Simple Syllabus PDFs",
            "2. Convert — Docling parses PDFs to markdown",
            "3. Strip boilerplate — repeated headers/footers/policies",
            "4. Normalize hierarchy — headings, sections, lists",
            "5. Recover images & tables — Gemini multimodal",
            "6. Validate — LLM + false-positive (bronze→silver→gold)",
            "7. Chunk & embed — Voyage voyage-3 → MongoDB Atlas",
        ],
    )

    add_section(
        s1,
        Inches(0.5),
        Inches(4.0),
        Inches(6.2),
        Inches(2.9),
        "Frontend & Ops",
        [
            "Streamlit chat UI w/ hot-reload for RAG modules",
            "Citation on every factual claim, linked to its source",
            "Inline links to the syllabus for each cited course",
            "Course comparison: tables with per-cell citations",
            "Docker + API proxy: TAMU/Voyage per-session budgets",
        ],
    )

    add_section(
        s1,
        Inches(6.9),
        Inches(4.0),
        Inches(6.0),
        Inches(2.9),
        "Observability & Eval",
        [
            "Langfuse per-node traces: cost, latency, tokens, errors",
            "Drill into any graph node for troubleshooting",
            "RAGAS answer quality: Faithfulness + Answer Relevancy",
            "RAGAS retrieval quality: Context Precision + Recall",
            "Golden-set benchmarks + chunking-eval harness",
            "pytest + mypy + ruff in CI loop",
        ],
    )

    add_footer(s1, "Branch: main  •  Recent: course comparison revamp, Docling v4, image_recovery filter")

    # Slide 3 — RAG Pipeline Flow (architecture diagram)
    s_flow = prs.slides.add_slide(blank)
    add_title_bar(s_flow, "RAG Pipeline — Flow")

    arch_path = Path(__file__).resolve().parents[1] / "docs" / "architecture.png"
    s_flow.shapes.add_picture(
        str(arch_path),
        Inches(0.4),
        Inches(1.1),
        width=Inches(8.6),
        height=Inches(5.9),
    )

    add_section(
        s_flow,
        Inches(9.2),
        Inches(1.1),
        Inches(3.9),
        Inches(5.9),
        "How a query flows",
        [
            "HISTORY_INJECT — pulls prior turns from mem0 Cloud",
            "ROUTER — TAMU LLM extracts variables, picks function",
            "Out-of-scope → CANNED_RESPONSE (early exit)",
            "RETRIEVAL — MongoDB Atlas (vector + BM25) + Voyage rerank",
            "GENERATOR — TAMU LLM produces cited answer",
            "Recursive loop — Generator can re-trigger Retrieval w/ updated state",
            "HISTORY_UPDATE — writes turn back to mem0",
            "Langfuse traces every node end-to-end",
        ],
    )

    add_footer(s_flow, "LangGraph state machine — Router → Retrieval → Generator with recursive refinement")

    # Slide 4 — Future Goals
    s2 = prs.slides.add_slide(blank)
    add_title_bar(s2, "TamuBot — Future Goals")

    add_section(
        s2,
        Inches(0.5),
        Inches(1.1),
        Inches(6.2),
        Inches(2.8),
        "Data Coverage",
        [
            "Add catalog.tamu.edu (course descriptions, degree plans, certificates)",
            "Switch from fixed-size to semantic chunking by syllabus category",
            "Close syllabus gaps across Howdy Portal + Simple Syllabus",
            "Expand beyond ISEN/CSCE to all graduate departments",
            "Promote silver → gold for the full corpus via iterative QA",
            "Refresh cadence: scheduled re-scrape per term",
        ],
    )

    add_section(
        s2,
        Inches(6.9),
        Inches(1.1),
        Inches(6.0),
        Inches(2.8),
        "Retrieval Quality",
        [
            "Per-stage LLMs (router/retrieval/generator) for speed + quality",
            "Optimize semantic search: embeddings, top-k, thresholds",
            "Multi-step / multi-hop query evaluation in the eval harness",
            "Schema-aware diffing for multi-course comparisons",
            "Recursive discovery: broaden anchor strategies",
            "Out-of-scope rejection: fewer false positives",
        ],
    )

    add_section(
        s2,
        Inches(0.5),
        Inches(4.0),
        Inches(6.2),
        Inches(2.9),
        "Product & UX",
        [
            "Polish frontend UX: layout, readability, response display",
            "Feedback loop: thumbs / ratings on each response",
            "Smoother, more reliable conversational flow",
            "Auth + persistent per-student memory across sessions",
            "Advisor mode: planning, prerequisites, degree-plan overlays",
            "Source preview pane: jump from citation → chunk",
        ],
    )

    add_section(
        s2,
        Inches(6.9),
        Inches(4.0),
        Inches(6.0),
        Inches(2.9),
        "Reliability & Deployment",
        [
            "Red-team the pipeline: prompt injection + data leakage tests",
            "Cost management: caching, budgets, quotas",
            "Load-balancing gateway in front of LLM calls",
            "Cloud deployment for production hosting",
            "Quality alerting: notify on regressions and drift",
            "Grow golden set; track regression on every PR",
        ],
    )

    add_footer(s2, "Targets shaped by current backlog and recent ingestion-pipeline iterations")

    # Slide 5 — Related Work: Similar Systems at Other Universities
    s3 = prs.slides.add_slide(blank)
    add_title_bar(s3, "Related Work — Similar Systems at Other Universities")

    # Left column: RAG-Based Academic Systems
    add_category_header(s3, Inches(0.5), Inches(1.05), Inches(6.2), "RAG-Based Academic Systems")

    add_system_card(
        s3, Inches(0.5), Inches(1.5), Inches(6.2), Inches(1.8),
        "Aurora", "Florida International University",
        "Neuro-symbolic advising agent — RAG + Prolog logic engine for policy-compliant course advising. 0.93 semantic alignment, 0.71s latency.",
        "arXiv:2602.17999",
    )
    add_system_card(
        s3, Inches(0.5), Inches(3.3), Inches(6.2), Inches(1.8),
        "BarkPlug V.2", "Mississippi State University",
        "RAG pipeline over 42 university resource documents using GPT-3.5-turbo. Evaluated with RAGAS (score 0.96) — same framework as TamuBot.",
        "arXiv:2405.08120",
    )
    add_system_card(
        s3, Inches(0.5), Inches(5.1), Inches(6.2), Inches(1.8),
        "URAG", "Ho Chi Minh City University of Technology",
        "Hybrid two-tier RAG (FAQ gate + vector retrieval) for admissions advising. Only live-deployed system of all six — running at ura.hcmut.edu.vn.",
        "arXiv:2501.16276",
    )

    # Right column: Broad AI Student Support Systems
    add_category_header(s3, Inches(6.9), Inches(1.05), Inches(6.0), "Broad AI Student Support Systems")

    add_system_card(
        s3, Inches(6.9), Inches(1.5), Inches(6.0), Inches(1.8),
        "Jill Watson", "Georgia Tech",
        "Oldest university AI assistant (2016); rebuilt on RAG + GPT-4o. Ingests syllabi and course Q&A; integrates via Canvas / Blackboard LTI.",
        "dilab.gatech.edu/jill-watson",
    )
    add_system_card(
        s3, Inches(6.9), Inches(3.3), Inches(6.0), Inches(1.8),
        "Maizey / Go Blue", "University of Michigan",
        "First institution-wide LLM deployment in U.S. higher ed (2023). 3,500+ RAG instances across courses, dining, events, campus services.",
        "its.umich.edu/computing/ai/go-blue-ai",
    )
    add_system_card(
        s3, Inches(6.9), Inches(5.1), Inches(6.0), Inches(1.8),
        "Smart Plan", "Purdue University",
        "AI-driven degree planning engine (\"GPS for graduation\"). Enforces prerequisites, models what-if major changes, integrates with SIS/registrar in real time.",
        "Purdue Registrar",
    )

    add_footer(s3, "Six representative academic RAG / AI systems — context for TamuBot's positioning")

    out = Path(__file__).resolve().parents[1] / "docs" / "tamubot_status.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    return out


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")
